"""Office/text extractors (D-04): DOCX, PPTX, XLSX, CSV, TXT.

python-docx / python-pptx / openpyxl for binary formats, stdlib csv for CSV
(rows joined with " | ") and plain UTF-8 read for TXT (errors="replace" so a
malformed byte never aborts extraction). openpyxl only reads the new .xlsx
format — legacy .xls is rejected with a clean ValueError (-> 400).
"""

import csv

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from app.documents.extractors import ExtractionResult


def extract_docx(path: str) -> str:
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs)


def extract_pptx(path: str) -> str:
    prs = Presentation(path)
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            # Only text-bearing shapes expose .text (python-pptx BaseShape
            # stubs don't declare it — getattr keeps pyright + runtime happy).
            text = getattr(shape, "text", None)
            if text:
                parts.append(str(text))
    return "\n".join(parts)


def extract_xlsx(path: str) -> str:
    wb = load_workbook(path, read_only=True, data_only=True)
    parts = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def extract_csv(path: str) -> str:
    with open(path, newline="", encoding="utf-8", errors="replace") as f:
        return "\n".join(" | ".join(row) for row in csv.reader(f))


def extract_txt(path: str) -> str:
    with open(path, encoding="utf-8", errors="replace") as f:
        return f.read()


async def extract_office(path: str, ext: str) -> ExtractionResult:
    if ext == ".docx":
        text = extract_docx(path)
    elif ext == ".pptx":
        text = extract_pptx(path)
    elif ext in {".xlsx", ".xls"}:
        try:
            text = extract_xlsx(path)
        except Exception as e:
            # openpyxl only reads the new .xlsx format — a legacy .xls raises
            # InvalidFileException (not a ValueError); convert for a clean 400.
            raise ValueError(f"unsupported or corrupt spreadsheet: {ext!r}") from e
    elif ext == ".csv":
        text = extract_csv(path)
    else:
        text = extract_txt(path)
    return ExtractionResult(text=text, pages=1, chars=len(text))
